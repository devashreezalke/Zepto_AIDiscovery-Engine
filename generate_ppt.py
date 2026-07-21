import sys
import subprocess
import os

# Auto-install python-pptx if it is not present
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not found. Installing now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Zepto Color Scheme
    purple = RGBColor(62, 0, 117)      # #3E0075
    magenta = RGBColor(226, 26, 132)   # #E21A84
    white = RGBColor(255, 255, 255)
    charcoal = RGBColor(55, 65, 81)     # #374151
    light_purple_bg = RGBColor(243, 234, 255) # #F3EAFF
    border_color = RGBColor(229, 231, 235)    # #E5E7EB

    # -------------------------------------------------------------
    # SLIDE 1: COGNITIVE HABIT LOOPS & EXPLORATION FRICTION
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Background color - off-white
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(249, 250, 251)
    
    # Slide Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cognitive Habit Loops and Trust Gaps Lock Users Into Narrow Grocery Checkout Cycles"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = purple
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "Analyzing why category exploration drops off over time and how platform discoverability fails to break routine loops."
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = charcoal
    
    # Content Columns (using shapes with background fills for professional aesthetic)
    # Left Column: Habit Loops
    left_shape = slide1.shapes.add_shape(
        1, Inches(0.75), Inches(2.0), Inches(5.7), Inches(4.7) # 1 = rectangle
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = white
    left_shape.line.color.rgb = border_color
    
    tf_left = left_shape.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = Inches(0.3)
    tf_left.margin_top = Inches(0.3)
    tf_left.margin_right = Inches(0.3)
    
    p1 = tf_left.paragraphs[0]
    p1.text = "THE HABIT LOOP LOCK-IN"
    p1.font.name = 'Arial'
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = purple
    
    bullets_left = [
        "User Inertia: Shopping behavior becomes highly repetitive to minimize cognitive load. Once a routine is established, users checkout in under 2 mins, bypassing catalog browsing.",
        "App Partitioning: Users mentally assign distinct roles to apps. Zepto is strictly anchored as 'immediate grocery staples,' while competitors (Nykaa, Amazon) capture adjacent categories.",
        "Checklist Focus: Users open the app with a specific checklist and exit instantly. Non-grocery tabs and discovery links are invisible to their active attention paths."
    ]
    for bullet in bullets_left:
        p_b = tf_left.add_paragraph()
        p_b.text = "• " + bullet
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = charcoal
        p_b.space_before = Pt(10)

    # Right Column: Exploration Barriers
    right_shape = slide1.shapes.add_shape(
        1, Inches(6.88), Inches(2.0), Inches(5.7), Inches(4.7)
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = white
    right_shape.line.color.rgb = border_color
    
    tf_right = right_shape.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = Inches(0.3)
    tf_right.margin_top = Inches(0.3)
    tf_right.margin_right = Inches(0.3)
    
    p2 = tf_right.paragraphs[0]
    p2.text = "PLATFORM TRUST & DISCOVERY BARRIERS"
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = magenta
    
    bullets_right = [
        "Discoverability? NO. Static banners and category list folders have a click-through rate (CTR) below 1.5%. Users ignore general promotional banners completely.",
        "Intuitive? PARTIALLY. Search indexing handles direct items well but fails to contextually recommend adjacent needs (e.g., diaper search does not show baby wipes).",
        "Trust Gaps: Quality skepticism prevents fresh produce adoption. Opaque return guidelines and scripted bots reject replacement requests, creating return anxiety on premium/electronic categories.",
        "Information Deficit: Omission of crowdsourced star ratings, feedback reviews, and sizing grids prevents users from entering unverified categories."
    ]
    for bullet in bullets_right:
        p_b = tf_right.add_paragraph()
        p_b.text = "• " + bullet
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = charcoal
        p_b.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 2: BUSINESS VALUE & KPI TREE
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Background color - off-white
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = RGBColor(249, 250, 251)
    
    # Slide Title Box
    title_box2 = slide2.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
    tf2 = title_box2.text_frame
    tf2.word_wrap = True
    p2_t = tf2.paragraphs[0]
    p2_t.text = "Breaking Purchase Silos Drives LTV: Mathematical KPI Tree for Cross-Category Trial"
    p2_t.font.name = 'Arial'
    p2_t.font.size = Pt(24)
    p2_t.font.bold = True
    p2_t.font.color.rgb = purple
    
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "Why cross-category trial matters for profit margins and the algebraic metrics tree to track Monthly Active Customer adoption."
    p2_sub.font.name = 'Arial'
    p2_sub.font.size = Pt(14)
    p2_sub.font.color.rgb = charcoal
    
    # Left Content Shape: Business Value
    val_shape = slide2.shapes.add_shape(
        1, Inches(0.75), Inches(2.0), Inches(4.5), Inches(4.7)
    )
    val_shape.fill.solid()
    val_shape.fill.fore_color.rgb = white
    val_shape.line.color.rgb = border_color
    
    tf_val = val_shape.text_frame
    tf_val.word_wrap = True
    tf_val.margin_left = Inches(0.3)
    tf_val.margin_top = Inches(0.3)
    tf_val.margin_right = Inches(0.3)
    
    p_v = tf_val.paragraphs[0]
    p_v.text = "WHY PRODUCT OUTCOME MATTERS"
    p_v.font.name = 'Arial'
    p_v.font.size = Pt(16)
    p_v.font.bold = True
    p_v.font.color.rgb = purple
    
    bullets_val = [
        "Margin Expansion: Staples (milk, eggs) have gross margins under 3-5%. Premium categories (Beauty, Pet, Electronics Accessories) yield margins of 18-25%.",
        "LTV & Frequency: Customers who adopt >= 3 categories have 3.4x higher Customer Lifetime Value and 80% lower churn than staples-only buyers.",
        "Basket Size Growth: Cross-category shoppers waive delivery limits easily, reducing delivery fee dependencies and shipping costs."
    ]
    for bullet in bullets_val:
        p_b = tf_val.add_paragraph()
        p_b.text = "• " + bullet
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = charcoal
        p_b.space_before = Pt(12)

    # Right Content Shape: KPI Tree
    kpi_shape = slide2.shapes.add_shape(
        1, Inches(5.68), Inches(2.0), Inches(6.9), Inches(4.7)
    )
    kpi_shape.fill.solid()
    kpi_shape.fill.fore_color.rgb = light_purple_bg
    kpi_shape.line.color.rgb = border_color
    
    tf_kpi = kpi_shape.text_frame
    tf_kpi.word_wrap = True
    tf_kpi.margin_left = Inches(0.4)
    tf_kpi.margin_top = Inches(0.3)
    tf_kpi.margin_right = Inches(0.4)
    
    p_k = tf_kpi.paragraphs[0]
    p_k.text = "MATHEMATICAL KPI TREE DESIGN"
    p_k.font.name = 'Arial'
    p_k.font.size = Pt(16)
    p_k.font.bold = True
    p_k.font.color.rgb = purple
    
    # KPI formulas with large font size
    formulas = [
        ("Primary North Star Metric (Monthly Category Trial Rate):", 
         "MCTR (%) = [ MACs purchasing from >= 1 new category in month t / Total MACs in month t ] x 100"),
        ("Sub-Metric Driver 1 (Discovery Click-Through Rate):", 
         "D-CTR (%) = [ Unique sessions clicking non-grocery PDPs / Total Sessions ] x 100"),
        ("Sub-Metric Driver 2 (Cross-Category Conversion Rate):", 
         "C-CONV (%) = [ Completed orders containing new category / Unique new PDP views ] x 100"),
        ("Actionable Product Funnel Drivers:", 
         "MCTR = f ( Search Redirect CTR + Cart Cross-sell % + Refund Badge trust conversion )")
    ]
    
    for title, formula in formulas:
        p_t = tf_kpi.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = charcoal
        p_t.space_before = Pt(8)
        
        p_f = tf_kpi.add_paragraph()
        p_f.text = "  " + formula
        p_f.font.name = 'Courier New'
        p_f.font.size = Pt(13.5)
        p_f.font.bold = True
        p_f.font.color.rgb = magenta
        p_f.space_before = Pt(2)
        
    # Save the presentation
    filename = "Zepto_Category_Discovery_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as '{os.path.abspath(filename)}'")

if __name__ == '__main__':
    create_presentation()
